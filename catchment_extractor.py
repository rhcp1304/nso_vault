import os
import io
import sys
import tempfile
import csv
import re  # Import the regular expression module
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from pptx import Presentation
from pathlib import Path

# --- Configuration ---
SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]
PARENT_FOLDER_ID = "1NuFnp9KFTmXmnluDdKKWp1J-r8dec1n_"
# Target the slide by its title instead of fixed index
SLIDE_TITLE = "Commercial Terms"
MIME_TYPE_PPTX = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
MIME_TYPE_FOLDER = 'application/vnd.google-apps.folder'

# List of all keys to extract
KEYS_TO_FIND = [
    "Catchment Name :",
    "Store Size",
    "Rent per Sq.ft",
    "Total Rent + Maintenance",
    "PROTO (in lakhs)",
    "GeoIQ Revenue Projection 2025 (in lakhs)"
]

# Define the output file name and header row for CSV
OUTPUT_FILE_NAME = "extracted_pptx_data.csv"
CSV_HEADERS = ["Source File Name"] + [key.replace(" :", "") for key in KEYS_TO_FIND]


# ---------------------

def get_drive_service():
    """Authenticates and returns the Google Drive API service object."""
    TOKEN_FILE = "token.json"
    creds = None
    if os.path.exists(TOKEN_FILE):
        creds = Credentials.from_authorized_user_file(TOKEN_FILE, SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            print("Authenticating with Google Drive...")
            try:
                flow = InstalledAppFlow.from_client_secrets_file("bdstorage_credentials.json", SCOPES)
                creds = flow.run_local_server(port=0)
            except FileNotFoundError:
                print("\n❌ ERROR: 'bdstorage_credentials.json' not found.")
                sys.exit(1)
        with open(TOKEN_FILE, "w") as token:
            token.write(creds.to_json())
    return build("drive", "v3", credentials=creds)


def extract_field_value(full_text, key):
    """
    Extracts the value following a specific key in a text block using regex
    for more robust key matching and value extraction.
    """
    # 1. Prepare key for regex: escape special characters, and remove trailing ':' if present
    # This key cleanup allows us to match keys like "Store Size" and "Store Size :"
    key_clean = re.escape(key.strip().rstrip(':')).replace(r'\ ', r'\s*')

    # 2. Construct the regex pattern:
    # (?:\s*:?\s*) - Allows for optional whitespace and an optional colon after the key
    # (.*?)         - Captures the value (non-greedy match)
    # (?:\n|\||$)   - Stops the capture at a newline (\n), a pipe (|), or the end of the string ($)
    pattern = rf"{key_clean}(?:\s*:?\s*)(.*?)(?:\n|\||$)"

    match = re.search(pattern, full_text, re.IGNORECASE | re.DOTALL)

    if match:
        try:
            # Group 1 holds the captured value
            final_value = match.group(1).strip()
            # If the value is empty, return None to signal failure to find data
            if not final_value:
                return None
            return final_value
        except Exception:
            return "[Extraction Error]"

    return None  # Key not found


def extract_data_from_ppt(file_path):
    """
    Reads a local .pptx file, finds the slide with the matching title, and extracts
    all key-value data points. Returns a dictionary of results.
    """
    results = {key: "" for key in KEYS_TO_FIND}
    keys_found_count = 0
    target_slide = None

    try:
        prs = Presentation(file_path)
    except Exception as e:
        results['error'] = f"[ERROR: Cannot open PPTX: {e}]"
        return results

    # --- 1. Find the target slide by title ---
    for i, slide in enumerate(prs.slides):
        slide_title = ""

        # Check title placeholder if available and non-empty
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            slide_title = slide.shapes.title.text.strip()

        # Fallback: check text of the first shape (in case a title placeholder isn't used)
        # We need to extract all text to find the title if it's not a placeholder
        elif slide.shapes and slide.shapes[0].has_text_frame:
            # Look for the first line of the first shape's text
            slide_title = slide.shapes[0].text_frame.text.strip().split('\n')[0]

        # Check for title match (case-insensitive)
        if slide_title and SLIDE_TITLE.lower() in slide_title.lower():
            target_slide = slide
            break

    if not target_slide:
        results['error'] = f"[ERROR: Slide with title '{SLIDE_TITLE}' not found]"
        return results

    # --- 2. Extract data from the found slide ---
    slide = target_slide

    # We need to iterate over all text-containing objects, including shapes AND tables

    # Process text from standard shapes (text boxes, headers)
    for shape in slide.shapes:
        # Check if shape is a text-containing object
        if shape.has_text_frame:
            full_text = shape.text_frame.text.strip()
            # This handles text boxes
            for key in KEYS_TO_FIND:
                if not results.get(key) or results[key] == "N/A":
                    value = extract_field_value(full_text, key)
                    if value is not None:
                        results[key] = value
                        keys_found_count += 1

        # Check if shape is a table
        elif shape.has_table:
            # This is the crucial part for table data
            table = shape.table

            # Iterate through rows and cells to extract all text from the table
            table_text = ""
            for row in table.rows:
                for cell in row.cells:
                    # Append all text from all cells to a single block, separated by newlines
                    # This allows the regex extractor to scan the entire table content
                    table_text += cell.text_frame.text.strip() + '\n'

            for key in KEYS_TO_FIND:
                if not results.get(key) or results[key] == "N/A":
                    value = extract_field_value(table_text, key)
                    if value is not None:
                        results[key] = value
                        keys_found_count += 1

    # If any key is still empty, mark it as 'N/A'
    for key in KEYS_TO_FIND:
        if not results.get(key):
            results[key] = "N/A"

    return results


def find_and_process_files(drive_service, current_folder_id, csv_writer):
    """
    Recursively finds all .pptx files and subfolders, processes files,
    and writes the extracted data to the CSV file.
    """
    # Query to list all files AND folders directly inside the current_folder_id
    query = (
        f"'{current_folder_id}' in parents and trashed=false"
    )

    page_token = None
    processed_count = 0

    while True:
        try:
            results = drive_service.files().list(
                q=query,
                spaces='drive',
                fields='nextPageToken, files(id, name, mimeType)',
                pageToken=page_token
            ).execute()
        except Exception as e:
            print(f"❌ Error listing contents of folder ID {current_folder_id}: {e}")
            break  # Stop processing this branch

        items = results.get('files', [])

        for item in items:
            item_id = item.get('id')
            item_name = item.get('name')
            item_mime = item.get('mimeType')

            # --- 1. Process PPTX Files ---
            if item_mime == MIME_TYPE_PPTX:
                print(f"|-- 📄 Processing: {item_name}")

                request = drive_service.files().get_media(fileId=item_id)
                file_handle = io.BytesIO()
                downloader = MediaIoBaseDownload(file_handle, request)

                done = False
                while not done:
                    status, done = downloader.next_chunk()
                file_handle.seek(0)

                with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
                    temp_file.write(file_handle.read())
                    temp_path = temp_file.name

                # Extract all structured data
                extracted_data = extract_data_from_ppt(temp_path)
                os.remove(temp_path)
                processed_count += 1

                # Prepare row for CSV
                row_data = [item_name]
                success = True

                if 'error' in extracted_data:
                    row_data.extend([extracted_data['error']] + [""] * (len(KEYS_TO_FIND) - 1))
                    success = False
                else:
                    # Append all extracted values in the order defined by KEYS_TO_FIND
                    for key in KEYS_TO_FIND:
                        row_data.append(extracted_data[key])
                        if extracted_data[key] == "N/A":
                            success = False

                # Write to CSV
                csv_writer.writerow(row_data)

                # Console feedback
                status_char = "✅" if success else "⚠️"
                print(f"|--   --> Status: {status_char} Catchment: {extracted_data.get(KEYS_TO_FIND[0], 'N/A')}")
                if not success:
                    print(f"|--   --> ERROR/MISSING DATA: {extracted_data.get('error', 'Some fields are N/A')}")


            # --- 2. Recurse into Subfolders ---
            elif item_mime == MIME_TYPE_FOLDER:
                print(f"|-- 📁 Found Folder: {item_name}. Entering...")
                # Recurse and accumulate the count
                processed_count += find_and_process_files(drive_service, item_id, csv_writer)
                print(f"|-- ⬆️ Exiting Folder: {item_name}")

        page_token = results.get('nextPageToken', None)
        if not page_token:
            break

    return processed_count


if __name__ == "__main__":

    if not Path("bdstorage_credentials.json").is_file():
        print("🔴 ERROR: 'bdstorage_credentials.json' not found.")
        print("Please ensure you have downloaded your Google Drive API credentials into the same folder.")
        sys.exit(1)

    print(
        f"\n--- Starting Recursive Scan in Folder ID: {PARENT_FOLDER_ID} (Looking for slide titled '{SLIDE_TITLE}') ---")

    try:
        drive_service = get_drive_service()
        total_processed_files = 0

        # Open file in 'w' mode (will overwrite) and set newline='' for proper CSV handling
        with open(OUTPUT_FILE_NAME, 'w', newline='', encoding='utf-8') as output_file:
            print(f"📝 Initializing output file: {OUTPUT_FILE_NAME}")

            # Initialize CSV writer
            csv_writer = csv.writer(output_file)

            # Write header row first
            csv_writer.writerow(CSV_HEADERS)

            # Start the recursive search and processing
            total_processed_files = find_and_process_files(drive_service, PARENT_FOLDER_ID, csv_writer)

        # --- Final Report ---
        if total_processed_files > 0:
            print("\n\n" + "=" * 50)
            print("✅ PROCESSING COMPLETE")
            print(f"Total PPTX files processed: {total_processed_files}")
            print(f"All results, including successful extractions and errors, are saved to: {OUTPUT_FILE_NAME}")
            print("==================================================")
        else:
            print("\n\n" + "=" * 50)
            print("No PPTX files found or processed in the Google Drive folder.")
            print("=" * 50)

    except Exception as e:
        print(f"\n❌ An unrecoverable error occurred: {e}")