import os
import re
import pickle
import io
import argparse
from google.auth.transport.requests import Request
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaFileUpload
from pptx import Presentation

# --- Configuration ---
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
CREDENTIALS_FILE = os.path.join(BASE_DIR, 'bdstorage_credentials.json')
TOKEN_FILE_PATH = os.path.join(BASE_DIR, 'token.pickle')

# --- Helper Functions ---

def authenticate_google_drive():
    SCOPES = ['https://www.googleapis.com/auth/drive']
    creds = None
    if os.path.exists(TOKEN_FILE_PATH):
        try:
            with open(TOKEN_FILE_PATH, 'rb') as token:
                creds = pickle.load(token)
            print("Loaded Drive API credentials from token file.")
        except Exception as e:
            print(f"Could not load Drive API token: {e}. Will re-authenticate.")
            creds = None

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            print("Drive API credentials expired, refreshing...")
            creds.refresh(Request())
        else:
            print(f"Initiating new Drive API authentication flow using {CREDENTIALS_FILE}...")
            if not os.path.exists(CREDENTIALS_FILE):
                raise FileNotFoundError(
                    f"Drive credentials file not found at {CREDENTIALS_FILE}. Please ensure it's there.")
            try:
                flow = InstalledAppFlow.from_client_secrets_file(
                    CREDENTIALS_FILE, SCOPES)
                creds = flow.run_local_server(port=0)
            except Exception as e:
                print(f"Error during OAuth flow: {e}")
                return None
        try:
            with open(TOKEN_FILE_PATH, 'wb') as token:
                pickle.dump(creds, token)
            print(f"Drive API credentials saved to {TOKEN_FILE_PATH}.")
        except Exception as e:
            print(f"Failed to save Drive API token: {e}")

    try:
        return build('drive', 'v3', credentials=creds)
    except Exception as e:
        print(f"Error building Drive service: {e}")
        return None

def get_market_and_zone_name_from_ppt(ppt_path):
    market_name = None
    zone_name = None
    try:
        prs = Presentation(ppt_path)
        if not prs.slides:
            print("PPT has no slides.")
            return None, None

        first_slide = prs.slides[0]
        slide_text = ""
        for shape in first_slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text += paragraph.text + "\n"

        zone_match = re.search(
            r"ZONE\s*:\s*(.*?)(?:\s*STATE|\s*CITY|\s*PIN CODE|$)",
            slide_text,
            re.IGNORECASE | re.DOTALL
        )
        if zone_match:
            zone_name = zone_match.group(1).strip()
            zone_name = re.sub(r'\s*\[Image \d+\]\s*', '', zone_name).strip()
        else:
            print("Could not find 'ZONE : ' on the first slide.")
            return None, None

        if zone_name:
            market_pattern = r"(?:^|\s)" + re.escape(zone_name) + r"\s+\d" + r"(?:.*?_){2,}.*?(?=\n|$)"
            market_match = re.search(
                market_pattern,
                slide_text,
                re.IGNORECASE | re.DOTALL
            )
            if market_match:
                market_name = market_match.group(0).strip()
                market_name = re.sub(r'\s*\[Image \d+\]\s*', '', market_name).strip()
                print(f"DEBUG: Found new market name: {market_name}")
            else:
                print(
                    f"Could not find a string starting with '{zone_name}' followed by a space and a digit, with at least two underscores, and ending at a newline.")

        if market_name is None:
            print(f"DEBUG: Extracted slide text:\n---START---\n{slide_text}\n---END---")

        return market_name, zone_name
    except Exception as e:
        print(f"An error occurred while reading the PPT: {e}")
        return None, None

def create_drive_folder(service, folder_name, parent_folder_id):
    file_metadata = {
        'name': folder_name,
        'mimeType': 'application/vnd.google-apps.folder',
        'parents': [parent_folder_id]
    }
    try:
        file = service.files().create(body=file_metadata, fields='id').execute()
        print(f"Folder '{folder_name}' created with ID: {file.get('id')}")
        return file.get('id')
    except HttpError as error:
        print(f"An HTTP error occurred during folder creation: {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred during folder creation: {e}")
        return None


def find_or_create_folder(service, folder_name, parent_folder_id):
    try:
        query = (
            f"name = '{folder_name}' and "
            f"mimeType = 'application/vnd.google-apps.folder' and "
            f"'{parent_folder_id}' in parents and "
            "trashed = false"
        )
        results = service.files().list(
            q=query,
            spaces='drive',
            fields='files(id, name)'
        ).execute()

        items = results.get('files', [])
        if items:
            print(f"Found existing folder '{folder_name}' with ID: {items[0]['id']}")
            return items[0]['id']
        else:
            print(f"Folder '{folder_name}' not found, creating it...")
            return create_drive_folder(service, folder_name, parent_folder_id)
    except HttpError as error:
        print(f"An HTTP error occurred while finding/creating folder '{folder_name}': {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred while finding/creating folder '{folder_name}': {e}")
        return None

# --- NEW HELPER FUNCTION FOR UPLOADING ---
def upload_file_to_drive(service, file_path, parent_folder_id):
    try:
        file_name = os.path.basename(file_path)
        file_metadata = {
            'name': file_name,
            'parents': [parent_folder_id]
        }
        media = MediaFileUpload(file_path, resumable=True)
        file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
        print(f"File '{file_name}' uploaded with ID: {file.get('id')}")
        return file.get('id')
    except HttpError as error:
        print(f"An HTTP error occurred during upload: {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred during upload: {e}")
        return None


# --- MAIN PROCESSING FUNCTION (Updated) ---
def main_processor(ppt_file_path, parent_folder_id):
    print("Starting PPT processing and folder creation using OAuth 2.0...")
    drive_service = authenticate_google_drive()
    if not drive_service:
        print("Authentication failed. Cannot proceed.")
        return

    # First, read the local file to get the names
    market_name, zone_name = get_market_and_zone_name_from_ppt(ppt_file_path)

    if not market_name:
        print("Could not extract market name. Folder not created and PPT not uploaded.")
        return

    print(f"Extracted Market Name: {market_name}")
    print(f"Extracted Zone Name: {zone_name}")

    # Determine the target parent folder for the market
    target_parent_for_market = parent_folder_id
    if zone_name:
        zone_folder_id = find_or_create_folder(drive_service, zone_name, parent_folder_id)
        if zone_folder_id:
            target_parent_for_market = zone_folder_id
        else:
            print(f"Failed to find or create Zone folder '{zone_name}'. Market folder will be created directly under the main parent.")

    # Create the final market folder
    market_folder_id = create_drive_folder(drive_service, market_name, target_parent_for_market)

    if not market_folder_id:
        print(f"Failed to create Market folder '{market_name}'. PPT file not uploaded.")
        return

    # Now upload the file directly into the newly created market folder
    uploaded_file_id = upload_file_to_drive(drive_service, ppt_file_path, market_folder_id)

    if uploaded_file_id:
        print("PPT file uploaded and placed in the correct folder.")
    else:
        print("Failed to upload PPT file to the drive.")
    print("Process completed successfully.")


# --- Main Execution Block (Updated) ---
if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description='Uploads a local PowerPoint file to Google Drive and organizes it into a new folder.'
    )
    parser.add_argument(
        '--ppt-file-path',
        type=str,
        help='Local file path of the PowerPoint presentation to upload.',
        required=True,
    )
    parser.add_argument(
        '--parent-folder-id',
        type=str,
        help='Google Drive File ID of the parent folder where the new market name folder will be created.',
        required=True,
    )

    args = parser.parse_args()

    if not os.path.exists(CREDENTIALS_FILE):
        raise FileNotFoundError(
            f"Error: OAuth client secrets file not found at '{CREDENTIALS_FILE}'. "
            "Please ensure 'bdstorage_credentials.json' is in the same directory as this script."
        )

    if not os.path.exists(args.ppt_file_path):
        print(f"Error: Local PPT file not found at '{args.ppt_file_path}'.")
    else:
        main_processor(args.ppt_file_path, args.parent_folder_id)