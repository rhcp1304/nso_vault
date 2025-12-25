import io
import re
import os
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from google.oauth2.credentials import Credentials
from pptx import Presentation

# PDF Generation Libraries
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors

# --- Configuration ---
ROOT_FOLDER_ID = '1PbLZ5yU2a2a4aacPiFpfF2uM7q0xGzkU'
CLIENT_SECRETS_FILE = 'bdstorage_credentials.json'
TOKEN_FILE = 'token.json'
OUTPUT_PDF = 'coordinates_list.pdf'
SCOPES = ['https://www.googleapis.com/auth/drive.readonly']


def get_drive_service():
    creds = None
    if os.path.exists(TOKEN_FILE):
        creds = Credentials.from_authorized_user_file(TOKEN_FILE, SCOPES)
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(CLIENT_SECRETS_FILE, SCOPES)
            creds = flow.run_local_server(port=0)
        with open(TOKEN_FILE, 'w') as token:
            token.write(creds.to_json())
    return build('drive', 'v3', credentials=creds)


def extract_coords(file_content):
    """Parses PPTX content to find LAT and LON values."""
    try:
        prs = Presentation(io.BytesIO(file_content))
        lat, lon = "N/A", "N/A"
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text
                    lat_match = re.search(r'LAT:\s*([\d.-]+)', text, re.IGNORECASE)
                    lon_match = re.search(r'LON:\s*([\d.-]+)', text, re.IGNORECASE)
                    if lat_match: lat = lat_match.group(1)
                    if lon_match: lon = lon_match.group(1)
            if lat != "N/A" and lon != "N/A":
                break
        return lat, lon
    except Exception:
        return "Error", "Error"


def process_folder(service, folder_id, data_list, styles):
    """Recursively crawls folders for PPTX files."""
    query = f"'{folder_id}' in parents and (mimeType = 'application/vnd.openxmlformats-officedocument.presentationml.presentation' or mimeType = 'application/vnd.google-apps.folder') and trashed = false"
    results = service.files().list(q=query, fields="files(id, name, mimeType)").execute()
    items = results.get('files', [])

    for item in items:
        if item['mimeType'] == 'application/vnd.google-apps.folder':
            process_folder(service, item['id'], data_list, styles)
        else:
            print(f"Processing: {item['name']}")
            try:
                request = service.files().get_media(fileId=item['id'])
                fh = io.BytesIO()
                downloader = MediaIoBaseDownload(fh, request)
                done = False
                while not done:
                    _, done = downloader.next_chunk()

                lat, lon = extract_coords(fh.getvalue())

                # Serial Number calculation: current length of data_list
                serial_no = len(data_list)

                # Use Paragraph for the filename to enable text wrapping
                file_name_p = Paragraph(item['name'], styles['BodyText'])

                # Append: [S.No, Property Name, Lat, Lon]
                data_list.append([serial_no, file_name_p, lat, lon])
            except Exception as e:
                print(f"Error downloading {item['name']}: {e}")


def generate_pdf(data):
    """Creates a PDF with a styled table including serial numbers and wrapping."""
    doc = SimpleDocTemplate(OUTPUT_PDF, pagesize=letter, leftMargin=30, rightMargin=30, topMargin=40, bottomMargin=40)
    elements = []
    styles = getSampleStyleSheet()

    # Title
    elements.append(Paragraph("Proposed Properties and Coordinates Report", styles['Title']))
    elements.append(Spacer(1, 12))

    #

    # Table configuration
    # Column widths: 40 for S.No, 310 for Name, 90 for Lat, 90 for Lon
    table = Table(data, colWidths=[40, 310, 90, 90], repeatRows=1)

    style = TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.darkblue),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (0, -1), 'CENTER'),  # Center align Serial Numbers
        ('ALIGN', (1, 0), (1, -1), 'LEFT'),  # Left align Property Names
        ('ALIGN', (2, 0), (-1, -1), 'CENTER'),  # Center align coordinates
        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 10),
        ('BACKGROUND', (0, 1), (-1, -1), colors.aliceblue),
    ])
    table.setStyle(style)

    elements.append(table)
    doc.build(elements)
    print(f"\nSuccessfully generated: {OUTPUT_PDF}")


if __name__ == '__main__':
    service = get_drive_service()
    report_styles = getSampleStyleSheet()

    # Initialize list with headers
    # Note: Column count must match data appended in process_folder
    report_data = [["S.No", "Property Name", "Latitude", "Longitude"]]

    print("Scanning Google Drive...")
    process_folder(service, ROOT_FOLDER_ID, report_data, report_styles)

    if len(report_data) > 1:
        generate_pdf(report_data)
    else:
        print("No files found.")