import os
import re
import pickle
import tempfile
from django.shortcuts import render
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_POST
from google.auth.transport.requests import Request
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaFileUpload
from pptx import Presentation

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CREDENTIALS_FILE = os.path.join(BASE_DIR, 'bdstorage_credentials.json')
TOKEN_FILE_PATH = os.path.join(BASE_DIR, 'token.pickle')


def authenticate_google_drive():
    SCOPES = ['https://www.googleapis.com/auth/drive']
    creds = None
    if os.path.exists(TOKEN_FILE_PATH):
        try:
            with open(TOKEN_FILE_PATH, 'rb') as token:
                creds = pickle.load(token)
            print("Loaded Drive API credentials from token file.")
        except Exception as e:
            print(f"Could not load Drive API token: {e}. Will re-authenticate.")
            creds = None

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            print("Drive API credentials expired, refreshing...")
            creds.refresh(Request())
        else:
            print(f"Initiating new Drive API authentication flow using {CREDENTIALS_FILE}...")
            if not os.path.exists(CREDENTIALS_FILE):
                raise FileNotFoundError(
                    f"Drive credentials file not found at {CREDENTIALS_FILE}. Please ensure it's there.")
            try:
                flow = InstalledAppFlow.from_client_secrets_file(
                    CREDENTIALS_FILE, SCOPES)
                creds = flow.run_local_server(port=0)
            except Exception as e:
                print(f"Error during OAuth flow: {e}")
                return None
        try:
            with open(TOKEN_FILE_PATH, 'wb') as token:
                pickle.dump(creds, token)
            print(f"Drive API credentials saved to {TOKEN_FILE_PATH}.")
        except Exception as e:
            print(f"Failed to save Drive API token: {e}")

    try:
        return build('drive', 'v3', credentials=creds)
    except Exception as e:
        print(f"Error building Drive service: {e}")
        return None


def get_market_and_zone_name_from_ppt(ppt_path):
    market_name = None
    zone_name = None
    try:
        prs = Presentation(ppt_path)
        if not prs.slides:
            print("PPT has no slides.")
            return None, None
        first_slide = prs.slides[0]
        slide_text = ""
        for shape in first_slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text += paragraph.text + "\n"
        zone_match = re.search(
            r"ZONE\s*:\s*(.*?)(?:\s*STATE|\s*CITY|\s*PIN CODE|$)",
            slide_text,
            re.IGNORECASE | re.DOTALL
        )
        if zone_match:
            zone_name = zone_match.group(1).strip()
            zone_name = re.sub(r'\s*\[Image \d+\]\s*', '', zone_name).strip()
        else:
            print("Could not find 'ZONE : ' on the first slide.")
            return None, None
        if zone_name:
            # New combined pattern to capture both types of market names
            market_pattern_combined = r"(?:^" + re.escape(zone_name) + r"\s*\d_.*?_.*$|^BD-.*$|^Add_.*$)"

            market_match = re.search(
                market_pattern_combined,
                slide_text,
                re.IGNORECASE | re.MULTILINE
            )
            if market_match:
                market_name = market_match.group(0).strip()
                market_name = re.sub(r'\s*\[Image \d+\]\s*', '', market_name).strip()
                print(f"DEBUG: Found new market name: {market_name}")
            else:
                print(
                    "Could not find a string starting with '{zone_name}' followed by a space and a digit, with at least two underscores, or a string starting with 'BD-' or 'Add_'.")
        if market_name is None:
            print(f"DEBUG: Extracted slide text:\n---START---\n{slide_text}\n---END---")
        return market_name, zone_name
    except Exception as e:
        print(f"An error occurred while reading the PPT: {e}")
        return None, None


def create_drive_folder(service, folder_name, parent_folder_id):
    file_metadata = {
        'name': folder_name,
        'mimeType': 'application/vnd.google-apps.folder',
        'parents': [parent_folder_id]
    }
    try:
        file = service.files().create(body=file_metadata, fields='id').execute()
        print(f"Folder '{folder_name}' created with ID: {file.get('id')}")
        return file.get('id')
    except HttpError as error:
        print(f"An HTTP error occurred during folder creation: {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred during folder creation: {e}")
        return None


def find_or_create_folder(service, folder_name, parent_folder_id):
    try:
        query = (
            f"name = '{folder_name}' and "
            f"mimeType = 'application/vnd.google-apps.folder' and "
            f"'{parent_folder_id}' in parents and "
            "trashed = false"
        )
        results = service.files().list(
            q=query,
            spaces='drive',
            fields='files(id, name)'
        ).execute()
        items = results.get('files', [])
        if items:
            print(f"Found existing folder '{folder_name}' with ID: {items[0]['id']}")
            return items[0]['id']
        else:
            print(f"Folder '{folder_name}' not found, creating it...")
            return create_drive_folder(service, folder_name, parent_folder_id)
    except HttpError as error:
        print(f"An HTTP error occurred while finding/creating folder '{folder_name}': {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred while finding/creating folder '{folder_name}': {e}")
        return None


def upload_file_to_drive(service, file_path, parent_folder_id):
    try:
        file_name = os.path.basename(file_path)
        file_metadata = {
            'name': file_name,
            'parents': [parent_folder_id]
        }
        media = MediaFileUpload(file_path, resumable=True)
        file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
        print(f"File '{file_name}' uploaded with ID: {file.get('id')}")
        return file.get('id')
    except HttpError as error:
        print(f"An HTTP error occurred during upload: {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred during upload: {e}")
        return None


def main_processor(ppt_file_path, parent_folder_id):
    print("Starting PPT processing and folder creation using OAuth 2.0...")
    drive_service = authenticate_google_drive()
    if not drive_service:
        return {'error': 'Authentication failed. Cannot proceed.'}

    market_name, zone_name = get_market_and_zone_name_from_ppt(ppt_file_path)
    if not market_name:
        return {'error': 'Could not extract market name. Folder not created and PPT not uploaded.'}

    print(f"Extracted Market Name: {market_name}")
    print(f"Extracted Zone Name: {zone_name}")

    target_parent_for_market = parent_folder_id
    if zone_name:
        zone_folder_id = find_or_create_folder(drive_service, zone_name, parent_folder_id)
        if zone_folder_id:
            target_parent_for_market = zone_folder_id
        else:
            print(
                f"Failed to find or create Zone folder '{zone_name}'. Market folder will be created directly under the main parent.")

    market_folder_id = find_or_create_folder(drive_service, market_name, target_parent_for_market)
    if not market_folder_id:
        return {'error': f"Failed to create Market folder '{market_name}'. PPT file not uploaded."}

    uploaded_file_id = upload_file_to_drive(drive_service, ppt_file_path, market_folder_id)
    if uploaded_file_id:
        print("PPT file uploaded and placed in the correct folder.")
        return {'message': 'File uploaded and organized successfully!', 'file_id': uploaded_file_id}
    else:
        print("Failed to upload PPT file to the drive.")
        return {'error': 'Failed to upload PPT file to the drive.'}


# --- Django Views ---
def upload_page(request):
    """
    Renders the HTML form for the file upload.
    """
    return render(request, 'uploader/index.html')

 ₹
@csrf_exempt
@require_POST
def process_upload(request):
    """
    Handles the POST request from the form, processes the file, and returns a JSON response.
    """
    try:
        uploaded_file = request.FILES.get('ppt_file')
        parent_folder_id = request.POST.get('parent_folder_id')

        if not uploaded_file or not parent_folder_id:
            return JsonResponse({'error': 'Missing file or parent folder ID.'}, status=400)

        # Create a temporary directory to store the uploaded file
        temp_dir = tempfile.mkdtemp()
        temp_file_path = os.path.join(temp_dir, uploaded_file.name)

        # Write the uploaded file content to the new path
        with open(temp_file_path, 'wb+') as f:
            for chunk in uploaded_file.chunks():
                f.write(chunk)

        # Run the main processor logic with the new file path
        result = main_processor(temp_file_path, parent_folder_id)

        # Clean up the temporary file and directory
        os.remove(temp_file_path)
        os.rmdir(temp_dir)

        if 'error' in result:
            return JsonResponse(result, status=500)

        return JsonResponse(result, status=200)

    except Exception as e:
        return JsonResponse({'error': str(e)}, status=500)