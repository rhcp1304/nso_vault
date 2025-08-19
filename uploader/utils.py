import os
import re
import pickle
import io
import tempfile
import shutil
from google.auth.transport.requests import Request
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaFileUpload, MediaIoBaseDownload
from pptx import Presentation
from pytubefix import YouTube, exceptions as pytube_exceptions
from google.oauth2.credentials import Credentials

# =================================================================================
# === Global Configuration and Pathing ===
# =================================================================================

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CREDENTIALS_FILE = os.path.join(BASE_DIR, 'bdstorage_credentials.json')
TOKEN_FILE_PATH = os.path.join(BASE_DIR, 'token.pickle')

# --- Shared Helper Functions ---
def authenticate_google_drive():
    SCOPES = ['https://www.googleapis.com/auth/drive']
    creds = None
    if os.path.exists(TOKEN_FILE_PATH):
        try:
            with open(TOKEN_FILE_PATH, 'rb') as token:
                creds = pickle.load(token)
            print("Loaded Drive API credentials from token file.")
        except Exception as e:
            print(f"Could not load Drive API token: {e}. Will re-authenticate.")
            creds = None
    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            print("Drive API credentials expired, refreshing...")
            creds.refresh(Request())
        else:
            print(f"Initiating new Drive API authentication flow using {CREDENTIALS_FILE}...")
            if not os.path.exists(CREDENTIALS_FILE):
                raise FileNotFoundError(
                    f"Drive credentials file not found at {CREDENTIALS_FILE}. Please ensure it's there.")
            try:
                flow = InstalledAppFlow.from_client_secrets_file(
                    CREDENTIALS_FILE, SCOPES)
                creds = flow.run_local_server(port=0)
            except Exception as e:
                print(f"Error during OAuth flow: {e}")
                return None
        try:
            with open(TOKEN_FILE_PATH, 'wb') as token:
                pickle.dump(creds, token)
            print(f"Drive API credentials saved to {TOKEN_FILE_PATH}.")
        except Exception as e:
            print(f"Failed to save Drive API token: {e}")
    try:
        return build('drive', 'v3', credentials=creds)
    except Exception as e:
        print(f"Error building Drive service: {e}")
        return None

def get_market_and_zone_name_from_ppt(ppt_path):
    market_name = None
    zone_name = None
    try:
        prs = Presentation(ppt_path)
        if not prs.slides:
            print("PPT has no slides.")
            return None, None
        first_slide = prs.slides[0]
        slide_text = ""
        for shape in first_slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text += paragraph.text + "\n"
        zone_match = re.search(
            r"ZONE\s*:\s*(.*?)(?:\s*STATE|\s*CITY|\s*PIN CODE|$)",
            slide_text,
            re.IGNORECASE | re.DOTALL
        )
        if zone_match:
            zone_name = zone_match.group(1).strip()
            zone_name = re.sub(r'\s*\[Image \d+\]\s*', '', zone_name).strip()
        else:
            print("Could not find 'ZONE : ' on the first slide.")
            return None, None
        if zone_name:
            market_pattern = r"^" + re.escape(zone_name) + r"\s*\d_.*?_.*$"
            market_match = re.search(
                market_pattern,
                slide_text,
                re.IGNORECASE | re.MULTILINE
            )
            if market_match:
                market_name = market_match.group(0).strip()
                market_name = re.sub(r'\s*\[Image \d+\]\s*', '', market_name).strip()
                print(f"DEBUG: Found new market name: {market_name}")
            else:
                print(
                    "Could not find a string starting with '{zone_name}' followed by a space and a digit, with at least two underscores.")
        if market_name is None:
            print(f"DEBUG: Extracted slide text:\n---START---\n{slide_text}\n---END---")
        return market_name, zone_name
    except Exception as e:
        print(f"An error occurred while reading the PPT: {e}")
        return None, None

def create_drive_folder(service, folder_name, parent_folder_id):
    file_metadata = {
        'name': folder_name,
        'mimeType': 'application/vnd.google-apps.folder',
        'parents': [parent_folder_id]
    }
    try:
        file = service.files().create(body=file_metadata, fields='id').execute()
        print(f"Folder '{folder_name}' created with ID: {file.get('id')}")
        return file.get('id')
    except HttpError as error:
        print(f"An HTTP error occurred during folder creation: {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred during folder creation: {e}")
        return None

def find_or_create_folder(service, folder_name, parent_folder_id):
    try:
        query = (
            f"name = '{folder_name}' and "
            f"mimeType = 'application/vnd.google-apps.folder' and "
            f"'{parent_folder_id}' in parents and "
            "trashed = false"
        )
        results = service.files().list(
            q=query,
            spaces='drive',
            fields='files(id, name)'
        ).execute()
        items = results.get('files', [])
        if items:
            print(f"Found existing folder '{folder_name}' with ID: {items[0]['id']}")
            return items[0]['id']
        else:
            print(f"Folder '{folder_name}' not found, creating it...")
            return create_drive_folder(service, folder_name, parent_folder_id)
    except HttpError as error:
        print(f"An HTTP error occurred while finding/creating folder '{folder_name}': {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred while finding/creating folder '{folder_name}': {e}")
        return None

def upload_file_to_drive(service, file_path, parent_folder_id):
    try:
        file_name = os.path.basename(file_path)
        file_metadata = {
            'name': file_name,
            'parents': [parent_folder_id]
        }
        media = MediaFileUpload(file_path, resumable=True)
        file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
        print(f"File '{file_name}' uploaded with ID: {file.get('id')}")
        return file.get('id')
    except HttpError as error:
        print(f"An HTTP error occurred during upload: {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred during upload: {e}")
        return None

def download_file_from_drive(service, file_id: str, destination_path: str):
    try:
        request = service.files().get_media(fileId=file_id)
        with open(destination_path, 'wb') as fh:
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while done is False:
                status, done = downloader.next_chunk()
        return True, "video/mp4"
    except HttpError as error:
        print(f"An error occurred during file download from Drive (ID: {file_id}): {error}")
        return False, None
    except Exception as e:
        print(f"An unexpected error occurred during file download (ID: {file_id}): {e}")
        return False, None

def download_youtube_video(youtube_url: str, destination_path: str):
    try:
        print(f"Downloading highest resolution video-only stream from: {youtube_url}")
        yt = YouTube(youtube_url)
        video_stream = yt.streams.filter(progressive=False, only_video=True).order_by('resolution').desc().first()
        if video_stream:
            print(f"Found video stream with resolution: {video_stream.resolution}")
            video_stream.download(output_path=os.path.dirname(destination_path),
                                  filename=os.path.basename(destination_path))
            return True, "video/mp4"
        else:
            print(f"No suitable video-only stream found for {youtube_url}")
            return False, None
    except pytube_exceptions.VideoUnavailable:
        print(f"Error: The YouTube video at '{youtube_url}' is unavailable.")
        return False, None
    except Exception as e:
        print(f"An error occurred while downloading YouTube video '{youtube_url}': {e}")
        return False, None

def check_file_exists(service, file_name: str, folder_id: str):
    query = f"name='{file_name}' and '{folder_id}' in parents and trashed=false"
    try:
        results = service.files().list(q=query, fields='files(id)').execute()
        items = results.get('files', [])
        return len(items) > 0
    except HttpError as error:
        print(f"An error occurred while checking for file existence: {error}")
        return False

def extract_all_potential_links_from_last_slide(pptx_file_path: str) -> list[dict]:
    if not os.path.exists(pptx_file_path) or not pptx_file_path.lower().endswith('.pptx'):
        print(f"Error: Invalid PPTX file path '{pptx_file_path}'")
        return []
    found_links_with_names = []
    url_pattern = re.compile(r'https?://[^\s\]\)\}>"]+|[\w.-]+\.[a-z]{2,}(?:/[^\s]*)?')
    try:
        prs = Presentation(pptx_file_path)
        if not prs.slides:
            print("No slides found in presentation.")
            return []
        last_slide = prs.slides[-1]
        print(f"Analyzing the last slide (Slide {len(prs.slides)}) for all potential links...")

        def get_text_from_cell(cell):
            if cell.text_frame:
                return " ".join("".join([run.text for run in p.runs]) for p in cell.text_frame.paragraphs).strip()
            return ""

        def find_urls_in_text_content(text_frame_obj, associated_name=None):
            for paragraph in text_frame_obj.paragraphs:
                full_text = "".join([run.text for run in paragraph.runs])
                for match in url_pattern.finditer(full_text):
                    found_links_with_names.append({'name': associated_name, 'link': match.group(0).strip()})
                for run in paragraph.runs:
                    if run.hyperlink.address:
                        found_links_with_names.append({'name': associated_name, 'link': run.hyperlink.address})

        for shape in last_slide.shapes:
            if hasattr(shape, 'action') and shape.action.hyperlink:
                found_links_with_names.append({'name': None, 'link': shape.action.hyperlink.address})
            if shape.has_text_frame:
                find_urls_in_text_content(shape.text_frame)
            if shape.has_table:
                table = shape.table
                name_col_idx = -1
                if table.rows:
                    header_row = table.rows[0]
                    for i, cell in enumerate(header_row.cells):
                        cell_text = get_text_from_cell(cell).lower().strip()
                        if "name" in cell_text or "store name" in cell_text:
                            name_col_idx = i
                            break
                for row_idx, row in enumerate(table.rows):
                    if row_idx == 0: continue
                    current_row_name = None
                    if name_col_idx != -1 and name_col_idx < len(row.cells):
                        current_row_name = get_text_from_cell(row.cells[name_col_idx])
                    for cell in row.cells:
                        if cell.text_frame:
                            find_urls_in_text_content(cell.text_frame, associated_name=current_row_name)
            if hasattr(shape, 'image') and hasattr(shape.image, 'hyperlink') and shape.image.hyperlink.address:
                found_links_with_names.append({'name': None, 'link': shape.image.hyperlink.address})
        unique_links_with_names = {}
        for item in found_links_with_names:
            link, name = item['link'], item['name']
            if link not in unique_links_with_names or (name and not unique_links_with_names[link]['name']):
                unique_links_with_names[link] = {'name': name, 'link': link}
        print(f"Successfully extracted {len(unique_links_with_names)} unique links from PPT.")
        return list(unique_links_with_names.values())
    except Exception as e:
        print(f"An error occurred while extracting links: {e}")
        return []

def main_processor(ppt_file_path, parent_folder_id):
    """
    Handles the PPT upload and folder creation.
    Returns the ID of the newly created folder.
    """
    print("Starting PPT processing and folder creation using OAuth 2.0...")
    drive_service = authenticate_google_drive()
    if not drive_service:
        return {'error': 'Authentication failed. Cannot proceed.'}

    market_name, zone_name = get_market_and_zone_name_from_ppt(ppt_file_path)
    if not market_name:
        return {'error': 'Could not extract market name. Folder not created and PPT not uploaded.'}

    print(f"Extracted Market Name: {market_name}")
    print(f"Extracted Zone Name: {zone_name}")

    target_parent_for_market = parent_folder_id
    if zone_name:
        zone_folder_id = find_or_create_folder(drive_service, zone_name, parent_folder_id)
        if zone_folder_id:
            target_parent_for_market = zone_folder_id
        else:
            print(
                f"Failed to find or create Zone folder '{zone_name}'. Market folder will be created directly under the main parent.")

    market_folder_id = find_or_create_folder(drive_service, market_name, target_parent_for_market)
    if not market_folder_id:
        return {'error': f"Failed to create Market folder '{market_name}'. PPT file not uploaded."}

    uploaded_file_id = upload_file_to_drive(drive_service, ppt_file_path, market_folder_id)
    if uploaded_file_id:
        print("PPT file uploaded and placed in the correct folder.")
        # Return the folder ID, not a generic success message
        return {'market_folder_id': market_folder_id, 'file_id': uploaded_file_id}
    else:
        print("Failed to upload PPT file to the drive.")
        return {'error': 'Failed to upload PPT file to the drive.'}
