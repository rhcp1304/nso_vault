import os
import re
import pickle
import tempfile
import base64
from datetime import datetime, timedelta
import shutil
import json
from io import BytesIO
import quopri  # 🛑 CRITICAL: Import for quoted-printable decoding

from django.shortcuts import render
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
from django.views.decorators.http import require_http_methods

from google.auth.transport.requests import Request
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaFileUpload, MediaIoBaseDownload
from pptx import Presentation

# --- Configuration ---
# BASE_DIR should point to the project root (e.g., nso_vault/)
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CREDENTIALS_FILE = os.path.join(BASE_DIR, 'bdstorage_credentials.json')
TOKEN_FILE_PATH = os.path.join(BASE_DIR, 'token.pickle')

SCOPES = [
    'https://www.googleapis.com/auth/drive',
    'https://www.googleapis.com/auth/gmail.readonly'
]

# Drive Link Pattern: Simplified to reliably find the ID after /d/
DRIVE_ID_PATTERN = re.compile(r'/d/([a-zA-Z0-9_-]+)')
DEFAULT_DRIVE_FILENAME = "Linked_Drive_Presentation.pptx"


# --- Google API Authentication ---

def authenticate_google_services():
    """Authenticates for both Drive and Gmail APIs using a single flow."""
    creds = None
    if os.path.exists(TOKEN_FILE_PATH):
        try:
            with open(TOKEN_FILE_PATH, 'rb') as token:
                creds = pickle.load(token)
            # print("Loaded API credentials from token file.")
        except Exception:
            creds = None

    if not creds or not creds.valid:
        if creds and creds.expired and creds.refresh_token:
            # print("Credentials expired, refreshing...")
            creds.refresh(Request())
        else:
            # print(f"Initiating new authentication flow using {CREDENTIALS_FILE}...")
            if not os.path.exists(CREDENTIALS_FILE):
                raise FileNotFoundError(
                    f"Credentials file not found at {CREDENTIALS_FILE}. Please ensure it's there.")
            try:
                flow = InstalledAppFlow.from_client_secrets_file(
                    CREDENTIALS_FILE, SCOPES)
                creds = flow.run_local_server(port=0)
            except Exception as e:
                print(f"Error during OAuth flow: {e}")
                return None, None
        try:
            with open(TOKEN_FILE_PATH, 'wb') as token:
                pickle.dump(creds, token)
            # print(f"API credentials saved to {TOKEN_FILE_PATH}.")
        except Exception as e:
            print(f"Failed to save API token: {e}")

    try:
        drive_service = build('drive', 'v3', credentials=creds)
        gmail_service = build('gmail', 'v1', credentials=creds)
        return drive_service, gmail_service
    except Exception as e:
        print(f"Error building API services: {e}")
        return None, None


# --- Gmail Specific Functions ---

def get_messages_with_ppt(service, user_id='me'):
    """
    Fetches messages using the absolute broadest query possible (date only).
    """
    date_two_days_ago = datetime.now() - timedelta(days=2)
    date_query = date_two_days_ago.strftime('%Y/%m/%d')

    # 🛑 CRITICAL FIX: Use the broadest query: date only.
    query = f'after:{date_query} AND NOT is:chat'
    # If this fails, the email is older than 48 hours.

    print(f"Searching Gmail with query: '{query}'")

    try:
        messages = []
        page_token = None
        while True:
            response = service.users().messages().list(
                userId=user_id,
                q=query,
                pageToken=page_token
            ).execute()

            messages.extend(response.get('messages', []))
            page_token = response.get('nextPageToken')
            if not page_token:
                break

        messages.reverse()
        return messages

    except HttpError as error:
        print(f'An HTTP error occurred while listing messages: {error}')
        return []
    except Exception as e:
        print(f'An unexpected error occurred: {e}')
        return []

def get_attachment_parts_recursively(parts):
    """
    Recursively searches all parts of a Gmail message for ANY part that has
    an attachmentId, accepting all physical attachments for processing.
    """
    all_file_parts = []

    if not parts:
        return all_file_parts

    for part in parts:
        # The definitive check for a downloadable physical attachment
        attachment_id = part.get('body', {}).get('attachmentId')

        # Accepting all attachments with an ID, letting the PPTX reader filter non-PPTs.
        if attachment_id:
            all_file_parts.append(part)

        # Recursively check nested parts
        if 'parts' in part:
            all_file_parts.extend(get_attachment_parts_recursively(part['parts']))

    return all_file_parts


def extract_drive_link(message_payload):
    """
    Final, aggressive, byte-level search for the Drive file ID.
    This bypasses most text decoding and regex issues.
    """
    if not message_payload:
        return None, None

    # Simple, non-greedy ID pattern to use in bytes
    DRIVE_ID_BYTE_PATTERN = re.compile(b'/d/([a-zA-Z0-9_-]+)')

    def recursive_link_search(parts):
        if not parts:
            return None, None

        for part in parts:
            encoding = next(
                (h.get('value') for h in part.get('headers', []) if h.get('name') == 'Content-Transfer-Encoding'),
                '').lower()
            data = part.get('body', {}).get('data')

            if data:
                try:
                    decoded_bytes = base64.urlsafe_b64decode(data.encode('UTF-8'))

                    if 'quoted-printable' in encoding:
                        # Un-encode the bytes if necessary
                        decoded_bytes = quopri.decodestring(decoded_bytes)

                    # 1. CRITICAL: Search for the ID directly in the raw decoded bytes
                    match_id_in_bytes = DRIVE_ID_BYTE_PATTERN.search(decoded_bytes)

                    if match_id_in_bytes:
                        file_id = match_id_in_bytes.group(1).decode('utf-8')

                        # 2. Extract Filename (soft attempt in text)
                        try:
                            # Use soft decoding just to try and find the filename
                            decoded_data = decoded_bytes.decode('utf-8', errors='ignore')
                            filename_match = re.search(r'([a-zA-Z0-9_ -]+\.pptx)', decoded_data, re.IGNORECASE)
                            filename = filename_match.group(1).strip() if filename_match else DEFAULT_DRIVE_FILENAME

                        except:
                            filename = DEFAULT_DRIVE_FILENAME

                        print(f"✅ FINAL BYTE-LEVEL FOUND: ID={file_id}, Filename={filename}")
                        return file_id, filename

                except Exception as e:
                    # print(f"Warning: Failed to process part data in byte search: {e}")
                    pass

            if 'parts' in part:
                file_id, filename = recursive_link_search(part['parts'])
                if file_id:
                    return file_id, filename

        return None, None

    return recursive_link_search(message_payload.get('parts', []))


def download_file_from_drive(drive_service, file_id, file_name, temp_dir_base):
    """
    Downloads a file from Google Drive based on its ID with robust error logging.
    """
    # 1. Ensure .pptx extension is present and clean the filename
    clean_file_name = file_name.strip()
    if not clean_file_name.lower().endswith('.pptx'):
        # Ensure only one .pptx extension exists
        clean_file_name = os.path.splitext(clean_file_name)[0].strip() + '.pptx'

    # 2. Use a safe filename for the local disk (removes characters that might cause OS issues)
    safe_file_name = re.sub(r'[^\w\-_\. ()]', '_', clean_file_name)

    try:
        # Request the media content
        request = drive_service.files().get_media(fileId=file_id)
        file_path = os.path.join(temp_dir_base, safe_file_name)

        # Download the file to a buffer (BytesIO)
        fh = BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False
        print(f"DIAGNOSTIC: Attempting Drive download for ID: {file_id}, Filename: {safe_file_name}")

        while done is False:
            status, done = downloader.next_chunk()

        # Write the buffer content to the local file
        with open(file_path, 'wb') as f:
            f.write(fh.getvalue())

        print(f"✅ Downloaded Drive link file successfully: {safe_file_name}.")
        return file_path

    except HttpError as error:
        # 🛑 CRITICAL: Capture the exact Drive API error
        error_details = error.content.decode('utf-8', errors='ignore')
        print("\n-------------------------------------------------------------")
        print(f'🛑 DRIVE API ERROR: The file was likely found but DOWNLOAD FAILED.')
        print(f'   ID: {file_id}')
        print(f'   Status: {error.resp.status} (403 usually means permission denied)')
        print(f'   Details: {error_details}')
        print("-------------------------------------------------------------\n")
        return None
    except Exception as e:
        print(f'An UNEXPECTED error occurred during Drive download (ID: {file_id}): {e}')
        return None


def download_attachment(drive_service, gmail_service, msg_id, user_id='me'):
    """
    Downloads ALL physical attachments OR the single linked Drive file.
    """
    downloaded_files = []
    message = None
    temp_dir_base = None

    try:
        temp_dir_base = tempfile.mkdtemp()

        # Attempt to retrieve the full message
        message = gmail_service.users().messages().get(userId=user_id, id=msg_id, format='full').execute()

        # CRITICAL CHECK: Ensure the message has a payload
        payload = message.get('payload')
        if not payload:
            print(f"WARNING: Message ID {msg_id} retrieved but has no payload. Skipping.")
            return downloaded_files, message

        # --- 1. DRIVE LINK CHECK (FIRST PRIORITY) ---
        drive_file_id, drive_file_name = extract_drive_link(payload)

        if drive_file_id:
            # Found a Drive link, download it directly using the Drive service
            file_path = download_file_from_drive(
                drive_service, drive_file_id, drive_file_name, temp_dir_base
            )
            if file_path:
                downloaded_files.append({
                    'path': file_path,
                    'filename': drive_file_name,
                    'temp_dir': temp_dir_base
                })
            # 🛑 CRITICAL: RETURN HERE to prevent physical attachment check on Drive-linked files
            return downloaded_files, message

        # --- 2. PHYSICAL ATTACHMENT CHECK (SECOND PRIORITY) ---

        # Search for all parts with an attachmentId
        all_file_parts = get_attachment_parts_recursively(payload.get('parts', []))

        for part in all_file_parts:
            attachment_id = part['body']['attachmentId']

            att_data = gmail_service.users().messages().attachments().get(
                userId=user_id, messageId=msg_id, id=attachment_id
            ).execute()

            data = att_data.get('data')
            file_data = base64.urlsafe_b64decode(data.encode('UTF-8'))

            # Use the filename provided by the API
            filename = part.get('filename')

            # Robust filename fallback and extension guarantee
            if not filename:
                filename = f"missing_name_{msg_id}_part{attachment_id}.pptx"
            elif not os.path.splitext(filename)[1]:
                filename += '.pptx'

            temp_file_path = os.path.join(temp_dir_base, filename)

            with open(temp_file_path, 'wb') as f:
                f.write(file_data)

            print(f"Downloaded attachment: {filename} (ID: {attachment_id}).")

            downloaded_files.append({
                'path': temp_file_path,
                'filename': filename,
                'temp_dir': temp_dir_base
            })

        return downloaded_files, message  # Final successful return

    except HttpError as error:
        # Diagnostic Printout
        print(f'API ERROR for ID {msg_id}. Status: {error.resp.status}. Details: {error.content.decode()}')
        return [], message

    except Exception as e:
        # Catches unexpected errors
        print(f'An UNEXPECTED error occurred during message processing for ID {msg_id}: {e}')
        return [], message

    finally:
        # 🛑 Make cleanup SAFE by wrapping it in its own try/except 🛑
        if not downloaded_files and temp_dir_base and os.path.exists(temp_dir_base):
            try:
                shutil.rmtree(temp_dir_base)
            except Exception as cleanup_e:
                # Log the cleanup failure but let the function proceed with its intended return
                print(f"WARNING: Failed to cleanup temp directory {temp_dir_base} for ID {msg_id}: {cleanup_e}")


# --- PPT Processing Functions ---

def get_market_and_zone_name_from_ppt(ppt_path):
    """
    Extracts market and zone names from the first slide of the PPT.
    """
    market_name = None
    zone_name = None
    slide_text = ""
    try:
        prs = Presentation(ppt_path)
        if not prs.slides:
            print("PPT has no slides.")
            return None, None
        first_slide = prs.slides[0]

        for shape in first_slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    slide_text += paragraph.text + "\n"

        # 1. Search for ZONE
        zone_match = re.search(
            r"ZONE\s*:\s*(.*?)(?:\s*STATE|\s*CITY|\s*PIN CODE|$)",
            slide_text,
            re.IGNORECASE | re.DOTALL
        )
        if zone_match:
            zone_name = zone_match.group(1).strip()
            zone_name = re.sub(r'\s*\[Image \d+\]\s*', '', zone_name).strip()
        else:
            print("EXTRACTION FAIL: Could not find 'ZONE : ' pattern on the first slide.")
            return None, None

        # 2. Search for MARKET
        if zone_name:
            market_pattern_combined = r"(?:^" + re.escape(zone_name) + r"\s*\d_.*?_.*$|^BD-.*$|^Add_.*$)"

            market_match = re.search(
                market_pattern_combined,
                slide_text,
                re.IGNORECASE | re.MULTILINE
            )
            if market_match:
                market_name = market_match.group(0).strip()
                market_name = re.sub(r'\s*\[Image \d+\]\s*', '', market_name).strip()
            else:
                print(
                    f"EXTRACTION FAIL: Found ZONE: '{zone_name}', but MARKET name did not match expected patterns.")

        return market_name, zone_name
    except Exception as e:
        print(f"An error occurred while reading the PPT: {e}")
        return None, None


# --- Drive Operations ---

def create_drive_folder(service, folder_name, parent_folder_id):
    """Creates a new folder."""
    file_metadata = {
        'name': folder_name,
        'mimeType': 'application/vnd.google-apps.folder',
        'parents': [parent_folder_id]
    }
    try:
        file = service.files().create(body=file_metadata, fields='id').execute()
        return file.get('id')
    except HttpError as error:
        print(f"An HTTP error occurred during folder creation: {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred during folder creation: {e}")
        return None


def find_or_create_folder(service, folder_name, parent_folder_id):
    """Finds an existing folder or creates a new one."""
    try:
        query = (
            f"name = '{folder_name}' and "
            f"mimeType = 'application/vnd.google-apps.folder' and "
            f"'{parent_folder_id}' in parents and "
            "trashed = false"
        )
        results = service.files().list(
            q=query,
            spaces='drive',
            fields='files(id, name)'
        ).execute()
        items = results.get('files', [])
        if items:
            return items[0]['id']
        else:
            return create_drive_folder(service, folder_name, parent_folder_id)
    except HttpError as error:
        print(f"An HTTP error occurred while finding/creating folder '{folder_name}': {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred while finding/creating folder '{folder_name}': {e}")
        return None


def upload_file_to_drive(service, file_path, parent_folder_id):
    """
    Uploads a file to Drive, replacing an existing file with the same name
    in the same folder if one is found.
    """
    file_name = os.path.basename(file_path)
    existing_file_id = None
    try:
        query = (
            f"name = '{file_name}' and "
            f"mimeType != 'application/vnd.google-apps.folder' and "
            f"'{parent_folder_id}' in parents and "
            "trashed = false"
        )
        results = service.files().list(
            q=query,
            spaces='drive',
            fields='files(id, name)'
        ).execute()

        items = results.get('files', [])
        if items:
            existing_file_id = items[0]['id']

    except HttpError as error:
        print(f"An HTTP error occurred while searching for existing file: {error}")
        existing_file_id = None

    try:
        # Use MediaFileUpload for efficient upload
        media = MediaFileUpload(file_path, resumable=True)

        if existing_file_id:
            file = service.files().update(
                fileId=existing_file_id,
                media_body=media,
                fields='id'
            ).execute()
        else:
            file_metadata = {
                'name': file_name,
                'parents': [parent_folder_id]
            }
            file = service.files().create(
                body=file_metadata,
                media_body=media,
                fields='id'
            ).execute()

        return file.get('id')

    except HttpError as error:
        print(f"An HTTP error occurred during upload/update: {error}")
        return None
    except Exception as e:
        print(f"An unexpected error occurred during upload/update: {e}")
        return None


# --- Main Processing Logic ---

def main_processor(ppt_file_path, parent_folder_id, drive_service):
    """
    Processes the PPT, determines the target folders (Zone/Market),
    and uploads the file, replacing existing ones if necessary.
    """
    # print("Starting PPT processing and folder organization...")

    market_name, zone_name = get_market_and_zone_name_from_ppt(ppt_file_path)
    if not market_name:
        return {'error': 'Could not extract market name. Folder not created and PPT not uploaded.'}

    # print(f"Extracted Market Name: {market_name}")
    # print(f"Extracted Zone Name: {zone_name}")

    target_parent_for_market = parent_folder_id
    if zone_name:
        zone_folder_id = find_or_create_folder(drive_service, zone_name, parent_folder_id)
        if zone_folder_id:
            target_parent_for_market = zone_folder_id
        else:
            print("Failed to find or create Zone folder. Market folder will be created directly under the main parent.")

    market_folder_id = find_or_create_folder(drive_service, market_name, target_parent_for_market)
    if not market_folder_id:
        return {'error': f"Failed to create Market folder '{market_name}'. PPT file not uploaded."}

    uploaded_file_id = upload_file_to_drive(drive_service, ppt_file_path, market_folder_id)
    if uploaded_file_id:
        # print("PPT file uploaded/replaced and placed in the correct folder.")
        return {'message': 'File uploaded and organized successfully!', 'file_id': uploaded_file_id}
    else:
        print("Failed to upload/replace PPT file to the drive.")
        return {'error': 'Failed to upload/replace PPT file to the drive.'}


# --- Django Views ---

def trigger_email_check_page(request):
    """
    Renders the HTML form to trigger the email check.
    """
    return render(request, 'drive_uploader/email_trigger.html')


@csrf_exempt
@require_http_methods(["POST"])
def check_email_and_upload(request):
    """
    Checks Gmail for ANY message with an attachment in the last 48 hours and
    attempts to process all identified attachments/linked Drive files as PPTs.
    """
    parent_folder_id = request.POST.get('parent_folder_id')

    if not parent_folder_id:
        return JsonResponse({'error': 'Missing parent folder ID.'}, status=400)

    # 1. Authenticate
    drive_service, gmail_service = authenticate_google_services()
    if not drive_service or not gmail_service:
        return JsonResponse({'error': 'Google API authentication failed. Check credentials/token.'}, status=500)

    # 2. Fetch messages
    messages = get_messages_with_ppt(gmail_service)

    if not messages:
        return JsonResponse({'message': 'No relevant emails found in the search window.'}, status=200)

    all_results = []
    # Use a set to track directories to clean, avoiding duplicates
    dirs_to_cleanup = set()

    # 3. Process each message
    for message in messages:
        msg_id = message['id']

        # ATTEMPT: Download all identifiable attachments OR Drive links
        downloaded_files, full_message_payload = download_attachment(drive_service, gmail_service, msg_id)

        # 4. Process all found/downloaded files
        if downloaded_files:
            # All downloaded files share the same temp directory
            temp_dir_to_clean = downloaded_files[0]['temp_dir']
            dirs_to_cleanup.add(temp_dir_to_clean)  # Add the base dir for cleanup

            for file_data in downloaded_files:
                temp_file_path = file_data['path']
                filename = file_data['filename']

                # Attempt to process ANY file, relying on the PPTX library to confirm the format
                try:
                    result = main_processor(temp_file_path, parent_folder_id, drive_service)

                    # If processing succeeds, the file was a valid PPTX/PPT
                    result['source_email_id'] = msg_id
                    result['source_filename'] = filename
                    all_results.append(result)
                except Exception as e:
                    # Log files that could not be read as a PPT/PPTX
                    result = {
                        'error': f'File skipped. Could not be read as PPTX/PPT. Exception: {e}',
                        'source_email_id': msg_id,
                        'source_filename': filename
                    }
                    all_results.append(result)
        else:
            all_results.append({
                'error': f'No downloadable attachment or Drive link found in message ID: {msg_id}',
                'source_email_id': msg_id
            })

    # 5. Clean up ALL temporary files and directories
    for temp_dir in dirs_to_cleanup:
        try:
            # We clean the directory regardless of whether the processing succeeded,
            # as the file was already uploaded to Drive.
            if os.path.exists(temp_dir):
                shutil.rmtree(temp_dir)

        except Exception as e:
            print(f"Warning: Failed to cleanup temp directory {temp_dir}: {e}")

    return JsonResponse({
                            'message': f'Email processing complete. {len(all_results)} files processed from {len(messages)} emails. Results attached.',
                            'results': all_results}, status=200)