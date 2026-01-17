# video_processor/services.py (Fully Fixed Code)

import os
import re
import pickle
import io  # Needed for MediaIoBaseDownload
import shutil  # New import for robust directory cleanup
from google.auth.transport.requests import Request
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from googleapiclient.http import MediaFileUpload, MediaIoBaseDownload
from pptx import Presentation
from pytubefix import YouTube, exceptions as pytube_exceptions

SCOPES_VIDEO = ['https://www.googleapis.com/auth/drive']
TOKEN_FILE_VIDEO = 'token.pickle'
CREDENTIALS_FILE_VIDEO = 'bdstorage_credentials.json'
API_SERVICE_NAME = 'drive'
API_VERSION = 'v3'


class Style:
    def SUCCESS(self, msg): return f"\033[92mSUCCESS: {msg}\033[0m"

    def ERROR(self, msg): return f"\033[91mERROR: {msg}\033[0m"

    def WARNING(self, msg): return f"\033[93mWARNING: {msg}\033[0m"

    def INFO(self, msg): return f"{msg}"


# =================================================================================
# === CORE FIX: FILENAME SANITIZATION FUNCTION ===
# =================================================================================

def sanitize_filename(filename):
    """
    Replaces characters that are illegal in file paths across various operating systems
    with an underscore to ensure safe file creation.
    """
    # Pattern to match characters generally considered unsafe/illegal for file paths:
    # /, \n, \r, \t, :, *, ?, ", <, >, |
    illegal_chars_pattern = r'[\\/:*?"<>|\n\r\t]'

    # Replace illegal characters with an underscore
    safe_filename = re.sub(illegal_chars_pattern, '_', filename).strip()

    # Ensure the filename is not empty after sanitization
    if not safe_filename:
        safe_filename = "sanitized_download"

    return safe_filename


class DriveHelper:
    def __init__(self):
        self.style = Style()

    def _log(self, message, style_func=None):
        if style_func:
            print(style_func(message))
        else:
            print(message)

    def get_authenticated_drive_service(self):
        # ... (Your existing get_authenticated_drive_service logic remains here) ...
        creds = None
        token_path = os.path.join(os.getcwd(), TOKEN_FILE_VIDEO)
        credentials_path = os.path.join(os.getcwd(), CREDENTIALS_FILE_VIDEO)

        if os.path.exists(token_path):
            self._log(f"Loading credentials from {TOKEN_FILE_VIDEO}...")
            with open(token_path, 'rb') as token:
                try:
                    creds = pickle.load(token)
                except Exception as e:
                    self._log(f"Error loading {TOKEN_FILE_VIDEO}: {e}. Re-authenticating...",
                              style_func=self.style.ERROR)

        if not creds or not creds.valid:
            if creds and creds.expired and creds.refresh_token:
                self._log("Refreshing expired credentials...")
                try:
                    creds.refresh(Request())
                except Exception as e:
                    self._log(f"Error refreshing token: {e}. Re-authenticating...", style_func=self.style.ERROR)
                    creds = None

            if not creds:
                self._log(f"No valid credentials found. Initiating authentication flow (check your browser)...")
                if not os.path.exists(credentials_path):
                    self._log(
                        f"ERROR: '{CREDENTIALS_FILE_VIDEO}' not found at '{credentials_path}'. Please ensure it's in your project root.",
                        style_func=self.style.ERROR)
                    return None
                flow = InstalledAppFlow.from_client_secrets_file(credentials_path, SCOPES_VIDEO)
                creds = flow.run_console()
            self._log(f"Saving new credentials to {TOKEN_FILE_VIDEO}...")
            with open(token_path, 'wb') as token:
                pickle.dump(creds, token)

        self._log("Authentication successful.", style_func=self.style.SUCCESS)
        return build(API_SERVICE_NAME, API_VERSION, credentials=creds)

    def find_pptx_in_drive_folder(self, service, folder_id: str):
        # ... (Your existing find_pptx_in_drive_folder logic remains here) ...
        self._log(f"Searching for a PPTX file in folder ID '{folder_id}'...")
        query = f"'{folder_id}' in parents and mimeType = 'application/vnd.openxmlformats-officedocument.presentationml.presentation' and trashed = false"
        try:
            results = service.files().list(q=query, spaces='drive', fields='files(id, name, mimeType)').execute()
            items = results.get('files', [])
            if not items:
                self._log(f"No PPTX file found in folder '{folder_id}'.", style_func=self.style.WARNING)
                return None, None
            else:
                if len(items) > 1:
                    self._log(
                        f"WARNING: Found {len(items)} PPTX files. Using the first one found: '{items[0]['name']}'.",
                        style_func=self.style.WARNING)
                return items[0]['id'], items[0]['name']
        except HttpError as error:
            self._log(f"An error occurred while searching for PPTX: {error}", style_func=self.style.ERROR)
            return None, None

    def download_file_from_drive(self, service, file_id: str, destination_path: str):
        # ... (Your existing download_file_from_drive logic remains here) ...
        try:
            request = service.files().get_media(fileId=file_id)
            # This open() call is what failed before, but now destination_path will be sanitized.
            with open(destination_path, 'wb') as fh:
                downloader = MediaIoBaseDownload(fh, request)
                done = False
                while done is False:
                    status, done = downloader.next_chunk()

            # Log successful download with the path's filename
            self._log(f"SUCCESS: File downloaded to: {os.path.basename(destination_path)}",
                      style_func=self.style.SUCCESS)
            return True, "video/mp4"

        except HttpError as error:
            self._log(f"An error occurred during file download from Drive (ID: {file_id}): {error}",
                      style_func=self.style.ERROR)
            return False, None
        except Exception as e:
            # If this error persists, the issue might be directory creation, but
            # it is most likely the filename issue you described, which is now fixed.
            self._log(f"An unexpected error occurred during file download (ID: {file_id}): {e}",
                      style_func=self.style.ERROR)
            return False, None

    def download_youtube_video(self, youtube_url: str, destination_path: str):
        # ... (Your existing download_youtube_video logic remains here) ...
        try:
            self._log(f"Downloading highest resolution video-only stream from: {youtube_url}")
            yt = YouTube(youtube_url)
            video_stream = yt.streams.filter(progressive=False, only_video=True).order_by('resolution').desc().first()
            if video_stream:
                self._log(f"Found video stream with resolution: {video_stream.resolution}")
                # We pass the SAFE destination_path here.
                video_stream.download(output_path=os.path.dirname(destination_path),
                                      filename=os.path.basename(destination_path))
                return True, "video/mp4"
            else:
                self._log(f"No suitable video-only stream found for {youtube_url}", style_func=self.style.ERROR)
                return False, None
        except pytube_exceptions.VideoUnavailable:
            self._log(f"Error: The YouTube video at '{youtube_url}' is unavailable.", style_func=self.style.ERROR)
            return False, None
        except Exception as e:
            self._log(f"An error occurred while downloading YouTube video '{youtube_url}': {e}",
                      style_func=self.style.ERROR)
            return False, None

    def upload_file_to_drive(self, service, file_name: str, file_path: str, mime_type: str, parent_folder_id: str):
        # ... (Your existing upload_file_to_drive logic remains here) ...
        file_metadata = {'name': file_name, 'parents': [parent_folder_id]}
        media = MediaFileUpload(file_path, mimetype=mime_type, resumable=True)
        try:
            file = service.files().create(body=file_metadata, media_body=media, fields='id').execute()
            return file.get('id')
        except HttpError as error:
            self._log(f"An error occurred during file upload ({file_name}): {error}", style_func=self.style.ERROR)
            return None
        except Exception as e:
            self._log(f"An unexpected error occurred during file upload ({file_name}): {e}",
                      style_func=self.style.ERROR)
            return None

    def extract_all_potential_links(self, pptx_file_path: str) -> list[dict]:
        """
        Iterates through ALL slides in the presentation to extract potential
        video links from shapes, text frames, and tables.
        """
        if not os.path.exists(pptx_file_path) or not pptx_file_path.lower().endswith('.pptx'):
            self._log(f"Error: Invalid PPTX file path '{pptx_file_path}'", style_func=self.style.ERROR)
            return []

        found_links_with_names = []
        url_pattern = re.compile(r'https?://[^\s\]\)\}>"]+')

        try:
            prs = Presentation(pptx_file_path)
            if not prs.slides:
                self._log("No slides found in presentation.")
                return []

            self._log(f"Analyzing {len(prs.slides)} slides for potential links...")

            def get_text_from_cell(cell):
                if cell.text_frame:
                    return " ".join("".join([run.text for run in p.runs]) for p in cell.text_frame.paragraphs).strip()
                return ""

            def find_urls_in_text_content(text_frame_obj, associated_name=None):
                for paragraph in text_frame_obj.paragraphs:
                    full_text = "".join([run.text for run in paragraph.runs])
                    # Find plain text matches
                    for match in url_pattern.finditer(full_text):
                        found_links_with_names.append({'name': associated_name, 'link': match.group(0).strip()})
                    # Find embedded hyperlinks
                    for run in paragraph.runs:
                        if run.hyperlink and run.hyperlink.address:
                            found_links_with_names.append({'name': associated_name, 'link': run.hyperlink.address})

            # Iterate through every slide
            for slide_index, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    # 1. Action Hyperlinks (e.g., on images or buttons)
                    if hasattr(shape, 'action') and shape.action.hyperlink and shape.action.hyperlink.address:
                        found_links_with_names.append({'name': None, 'link': shape.action.hyperlink.address})

                    # 2. Text Frames
                    if shape.has_text_frame:
                        find_urls_in_text_content(shape.text_frame)

                    # 3. Tables (Common for lists of links/stores)
                    if shape.has_table:
                        table = shape.table
                        name_col_idx = -1
                        if table.rows:
                            header_row = table.rows[0]
                            for i, cell in enumerate(header_row.cells):
                                cell_text = get_text_from_cell(cell).lower().strip()
                                if any(x in cell_text for x in ["name", "store", "market"]):
                                    name_col_idx = i
                                    break

                        for row_idx, row in enumerate(table.rows):
                            # Optionally skip header, but check all cells for URLs
                            current_row_name = None
                            if name_col_idx != -1 and name_col_idx < len(row.cells):
                                current_row_name = get_text_from_cell(row.cells[name_col_idx])

                            for cell in row.cells:
                                if cell.text_frame:
                                    find_urls_in_text_content(cell.text_frame, associated_name=current_row_name)

            # Deduplicate results
            unique_links_with_names = {}
            for item in found_links_with_names:
                link, name = item['link'], item['name']
                if link not in unique_links_with_names or (name and not unique_links_with_names[link]['name']):
                    unique_links_with_names[link] = {'name': name, 'link': link}

            return list(unique_links_with_names.values())

        except Exception as e:
            self._log(f"An error occurred while extracting links: {e}", style_func=self.style.ERROR)
            return []

    def get_market_name_prefix(self, pptx_file_path: str) -> str:
        # ... (Your existing get_market_name_prefix logic remains here) ...
        """
        Extracts the market name from the first slide of a PPTX file.
        The market name is identified by a pattern: 'ZONE : [zone name]' followed by
        a line starting with the zone name, a digit, and two underscores.
        """
        if not os.path.exists(pptx_file_path):
            self._log(f"Error: PPTX file not found locally at '{pptx_file_path}'", style_func=self.style.ERROR)
            return ""
        try:
            prs = Presentation(pptx_file_path)
            if not prs.slides:
                self._log("PPT has no slides.", style_func=self.style.WARNING)
                return ""
            first_slide = prs.slides[0]
            slide_text = ""
            for shape in first_slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_text += paragraph.text + "\n"

            # 1. Search for the ZONE name
            zone_match = re.search(
                r"ZONE\s*:\s*(.*?)(?:\s*STATE|\s*CITY|\s*PIN CODE|$)",
                slide_text,
                re.IGNORECASE | re.DOTALL
            )
            if not zone_match:
                self._log("Could not find 'ZONE : ' on the first slide.", style_func=self.style.WARNING)
                return ""

            zone_name = zone_match.group(1).strip()
            zone_name = re.sub(r'\s*\[Image \d+\]\s*', '', zone_name).strip()

            if not zone_name:
                self._log("Found 'ZONE : ' but the zone name was empty.", style_func=self.style.WARNING)
                return ""

            # 2. Use the zone name to find the market name
            market_pattern = r"^" + re.escape(zone_name) + r"\s*\d_.*?_.*$"
            market_match = re.search(
                market_pattern,
                slide_text,
                re.IGNORECASE | re.MULTILINE
            )

            if market_match:
                full_market_name = market_match.group(0).strip()
                full_market_name = re.sub(r'\s*\[Image \d+\]\s*', '', full_market_name).strip()

                # Apply the original logic to extract the prefix from the found market name
                if '_' in full_market_name:
                    prefix = full_market_name.rsplit('_', 1)[1].strip()
                    self._log(f"Extracted market name prefix: '{prefix}'", style_func=self.style.INFO)
                    return prefix
                else:
                    self._log(f"No underscore found in market name. Using full value: '{full_market_name}'",
                              style_func=self.style.WARNING)
                    return full_market_name
            else:
                self._log(
                    f"Could not find a string starting with '{zone_name}' followed by a digit and two underscores.",
                    style_func=self.style.WARNING)
                return ""
        except Exception as e:
            self._log(f"An error occurred while extracting market name prefix: {e}", style_func=self.style.ERROR)
            return ""

    def check_file_exists(self, service, file_name: str, folder_id: str):
        # ... (Your existing check_file_exists logic remains here) ...
        query = f"name='{file_name}' and '{folder_id}' in parents and trashed=false"
        try:
            results = service.files().list(q=query, fields='files(id)').execute()
            items = results.get('files', [])
            return len(items) > 0
        except HttpError as error:
            self._log(f"An error occurred while checking for file existence: {error}", style_func=self.style.ERROR)
            return False


# --- Main Processing Function ---
def process_video_links_internal(google_drive_folder_id, temp_download_dir):
    drive_helper = DriveHelper()
    style = drive_helper.style
    print(style.SUCCESS(f"Starting process for Google Drive folder ID: {google_drive_folder_id}"))

    # Initialize service and folder check
    service = drive_helper.get_authenticated_drive_service()
    if not service:
        print(style.ERROR("Failed to authenticate with Google Drive API."))
        return

    # 1. Setup temporary directory (using a single try/finally cleanup block)
    try:
        if os.path.exists(temp_download_dir):
            shutil.rmtree(temp_download_dir)  # Clean up any previous run's failed state
            print(f"Cleaned previous temporary directory: '{temp_download_dir}'")

        os.makedirs(temp_download_dir)
        print(f"Created temporary download directory: '{temp_download_dir}'")

        # 2. Find PPTX
        pptx_drive_id, pptx_file_name_original = drive_helper.find_pptx_in_drive_folder(service, google_drive_folder_id)
        if not pptx_drive_id:
            print(style.ERROR(f"No PPTX file found in folder '{google_drive_folder_id}'."))
            return

        # 3. Download PPTX
        # ⭐️ FIX: SANITIZE THE FILENAME BEFORE CREATING THE LOCAL PATH ⭐️
        pptx_file_name_safe = sanitize_filename(pptx_file_name_original)
        local_pptx_path = os.path.join(temp_download_dir, pptx_file_name_safe)

        print(f"Downloading PPTX '{pptx_file_name_original}' (ID: {pptx_drive_id})...")
        if not drive_helper.download_file_from_drive(service, pptx_drive_id, local_pptx_path)[0]:
            print(style.ERROR(f"Failed to download PPTX: {pptx_file_name_original}"))
            return

        # 4. Extract links and prefixes
        market_name_prefix_raw = drive_helper.get_market_name_prefix(local_pptx_path)

        # Sanitize prefix for the final video file name on Drive
        cleaned_name = re.sub(r'[\\/:*?\"<>|]', '', market_name_prefix_raw).strip()
        prefix_for_filename = f"{cleaned_name} " if market_name_prefix_raw else ""

        print(
            f"Using market name prefix: '{prefix_for_filename}'" if prefix_for_filename else "No valid market name prefix found.")
        print("Extracting potential video links and associated names from PPTX...")

        extracted_links_with_names = drive_helper.extract_all_potential_links(local_pptx_path)
        print(f"Found {len(extracted_links_with_names)} potential links.")

        google_drive_file_id_pattern = re.compile(r'drive\.google\.com/(?:file/d/|uc\?id=)([a-zA-Z0-9_-]+)')
        youtube_pattern = re.compile(r'(?:youtube\.com/watch\?v=|youtu\.be/)([\w-]+)')
        processed_links = set()

        # 5. Process Videos
        for item in extracted_links_with_names:
            link, suggested_name = item['link'], item['name']
            if link in processed_links:
                continue
            processed_links.add(link)
            drive_match = google_drive_file_id_pattern.search(link)
            youtube_match = youtube_pattern.search(link)
            video_mime_type = "video/mp4"
            video_downloaded = False
            local_video_path = None

            if drive_match:
                video_drive_id = drive_match.group(1)
                print(f"\n--- Detected Google Drive video link: {link} (ID: {video_drive_id}) ---")
                try:
                    file_metadata = service.files().get(fileId=video_drive_id, fields='name,mimeType').execute()
                    original_video_name_from_drive = file_metadata.get('name', f"unknown_video_{video_drive_id}")

                    # Sanitize base name for the file name stored on Drive
                    base_name_for_file = re.sub(r'[\\/:*?"<>|]', '', (
                            suggested_name or os.path.splitext(original_video_name_from_drive)[0])).strip()
                    ext_from_drive = os.path.splitext(original_video_name_from_drive)[1]
                    final_video_name_for_drive = f"{prefix_for_filename}{base_name_for_file}{ext_from_drive}"

                    if drive_helper.check_file_exists(service, final_video_name_for_drive, google_drive_folder_id):
                        print(style.WARNING(
                            f"File '{final_video_name_for_drive}' already exists. Skipping download and upload."))
                        continue

                    # ⭐️ FIX: SANITIZE THE FILENAME FOR LOCAL DOWNLOAD ⭐️
                    local_video_name_safe = sanitize_filename(final_video_name_for_drive)
                    local_video_path = os.path.join(temp_download_dir, local_video_name_safe)

                    print(f"Downloading '{final_video_name_for_drive}'...")
                    video_downloaded, _ = drive_helper.download_file_from_drive(service, video_drive_id,
                                                                                local_video_path)

                except HttpError as api_error:
                    print(style.ERROR(f"Drive API error for link {link}: {api_error}"))
                    continue

            elif youtube_match:
                print(f"\n--- Detected YouTube video link: {link} ---")

                base_name = re.sub(r'[\\/:*?"<>|]', '', (suggested_name or 'youtube_video')).strip()
                final_video_name_for_drive = f"{prefix_for_filename}{base_name}.mp4"

                if drive_helper.check_file_exists(service, final_video_name_for_drive, google_drive_folder_id):
                    print(style.WARNING(
                        f"File '{final_video_name_for_drive}' already exists. Skipping download and upload."))
                    continue

                # ⭐️ FIX: SANITIZE THE FILENAME FOR LOCAL DOWNLOAD ⭐️
                local_video_name_safe = sanitize_filename(final_video_name_for_drive)
                local_video_path = os.path.join(temp_download_dir, local_video_name_safe)

                print(f"Downloading '{final_video_name_for_drive}'...")
                video_downloaded, _ = drive_helper.download_youtube_video(link, local_video_path)

            else:
                print(f"Skipping unsupported link: {link}")
                continue

            if video_downloaded and local_video_path and os.path.exists(local_video_path):

                if os.path.getsize(local_video_path) < 1024:
                    print(style.WARNING(f"WARNING: Downloaded video is very small. Skipping upload and cleaning up."))
                    os.remove(local_video_path)
                    continue

                print(f"Uploading '{final_video_name_for_drive}' to Google Drive...")
                uploaded_file_id = drive_helper.upload_file_to_drive(
                    service,
                    final_video_name_for_drive,
                    local_video_path,
                    video_mime_type,
                    google_drive_folder_id
                )

                if uploaded_file_id:
                    print(
                        style.SUCCESS(f"Successfully uploaded: {final_video_name_for_drive} (ID: {uploaded_file_id})"))
                else:
                    print(style.ERROR(f"Failed to upload '{final_video_name_for_drive}'."))

            elif not video_downloaded:
                print(style.ERROR(f"Failed to download video from {link}"))

    except Exception as e:
        print(style.ERROR(f"An unrecoverable error occurred during processing: {e}"))
        raise  # Re-raise the exception for Celery to handle

    finally:
        # 6. Final cleanup (Robust cleanup)
        if os.path.exists(temp_download_dir):
            try:
                shutil.rmtree(temp_download_dir)
                print(f"Cleaned up temporary directory: {temp_download_dir}")
            except OSError as e:
                print(style.WARNING(f"Could not remove directory '{temp_download_dir}': {e}"))

    print(style.SUCCESS("Process completed."))
    return True

# Add this to the bottom of video_processor/services.py

# REPLACE the existing run_recursive_video_automation at the end of services.py with this:

def run_recursive_video_automation(root_folder_id, base_temp_dir):
    """
    Scans Google Drive starting from root_folder_id.
    If a folder contains a PPTX, it triggers the download logic.
    """
    helper = DriveHelper()

    # FIX 1: Use the actual method name you defined in your class
    service = helper.get_authenticated_drive_service()

    style = helper.style

    if not service:
        print(style.ERROR("Failed to authenticate Drive service for recursion."))
        return

    def crawl(current_folder_id, current_folder_name):
        print(f"\n{style.INFO('=' * 60)}")
        print(f"📂 SCANNING FOLDER: {current_folder_name} (ID: {current_folder_id})")
        print(f"{style.INFO('=' * 60)}")

        # 1. Check if this specific folder has a PPTX
        query = f"'{current_folder_id}' in parents and mimeType = 'application/vnd.openxmlformats-officedocument.presentationml.presentation' and trashed = false"
        try:
            results = service.files().list(q=query, fields='files(id, name)').execute()
            pptx_files = results.get('files', [])

            if pptx_files:
                print(style.SUCCESS(f"Found PPTX in '{current_folder_name}'. Triggering download logic..."))

                # Create a unique temp path for this specific folder
                folder_temp_path = os.path.join(base_temp_dir, current_folder_id)

                # CALL YOUR EXISTING FUNCTION (no changes to your main logic)
                process_video_links_internal(current_folder_id, folder_temp_path)
            else:
                print(f"No PPTX in '{current_folder_name}', checking subfolders...")

            # 2. Find all subfolders inside this folder to continue crawling
            sub_query = f"'{current_folder_id}' in parents and mimeType = 'application/vnd.google-apps.folder' and trashed = false"
            sub_results = service.files().list(q=sub_query, fields='files(id, name)').execute()
            subfolders = sub_results.get('files', [])

            for folder in subfolders:
                crawl(folder['id'], folder['name'])

        except Exception as e:
            print(style.ERROR(f"Error crawling folder {current_folder_name}: {e}"))

    # Start the recursive loop
    crawl(root_folder_id, "ROOT")

    # FIX 2: Removed the systemctl stop command so the worker stays Always-On.
    print(style.SUCCESS("\n✅ RECURSIVE AUTOMATION COMPLETE. Worker remains IDLE for next task."))